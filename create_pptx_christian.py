#!/usr/bin/env python3
"""
Script para generar presentación PPTX para reunión con Christian
Estilo: Dark Premium con colores Emerald, Blue, Violet
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colores del tema
DARK_BG = RgbColor(0x12, 0x12, 0x1A)  # #12121A
DARK_BG_2 = RgbColor(0x0A, 0x0A, 0x12)  # Más oscuro para sombras
EMERALD = RgbColor(0x22, 0xC5, 0x5E)  # #22C55E
EMERALD_DARK = RgbColor(0x16, 0x8A, 0x42)  # Emerald oscuro
EMERALD_GLOW = RgbColor(0x34, 0xD3, 0x99)  # Emerald brillante
BLUE = RgbColor(0x3B, 0x82, 0xF6)     # #3B82F6
BLUE_DARK = RgbColor(0x1E, 0x40, 0xAF)  # Blue oscuro
BLUE_GLOW = RgbColor(0x60, 0xA5, 0xFA)  # Blue brillante
VIOLET = RgbColor(0x8B, 0x5C, 0xF6)   # #8B5CF6
VIOLET_DARK = RgbColor(0x6D, 0x28, 0xD9)  # Violet oscuro
VIOLET_GLOW = RgbColor(0xA7, 0x8B, 0xFA)  # Violet brillante
CYAN = RgbColor(0x06, 0xB6, 0xD4)     # Cyan accent
PINK = RgbColor(0xEC, 0x48, 0x99)     # Pink accent
WHITE = RgbColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RgbColor(0xA0, 0xA0, 0xA0)
GRAY_MEDIUM = RgbColor(0x6B, 0x6B, 0x80)
GRAY_DARK = RgbColor(0x2A, 0x2A, 0x35)
GLASS_BG = RgbColor(0x1E, 0x1E, 0x2E)  # Glass card background
GLASS_BORDER = RgbColor(0x3D, 0x3D, 0x50)  # Glass border

def set_slide_background(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Arial"):
    """Add a text box to slide"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rounded_rectangle(slide, left, top, width, height, fill_color,
                          transparency=0, line_color=None):
    """Add a rounded rectangle shape"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_color):
    """Add a circle shape"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_arrow(slide, left, top, width, height, fill_color):
    """Add an arrow shape"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_glass_card(slide, left, top, width, height, accent_color, title="", icon=""):
    """Create a glassmorphism card with glow effect and accent border"""
    # Shadow layer (offset behind)
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + Inches(0.08), top + Inches(0.08), width, height
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = DARK_BG_2
    shadow.line.fill.background()

    # Glow layer (subtle colored glow)
    glow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left - Inches(0.02), top - Inches(0.02),
        width + Inches(0.04), height + Inches(0.04)
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = accent_color
    glow.line.fill.background()

    # Main glass card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = GLASS_BG
    card.line.color.rgb = GLASS_BORDER
    card.line.width = Pt(1)

    # Top accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + Inches(0.1), top + Inches(0.1),
        width - Inches(0.2), Inches(0.06)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent_color
    accent_bar.line.fill.background()

    # Title with icon
    if icon or title:
        add_text_box(slide, left + Inches(0.15), top + Inches(0.25),
                     width - Inches(0.3), Inches(0.4),
                     f"{icon}  {title}" if icon else title,
                     font_size=14, font_color=accent_color, bold=True)

    return card

def add_mini_card(slide, left, top, width, height, bg_color=None):
    """Create a mini inner card"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color or RgbColor(0x16, 0x16, 0x22)
    card.line.color.rgb = RgbColor(0x2D, 0x2D, 0x3D)
    card.line.width = Pt(0.5)
    return card

def add_badge(slide, left, top, width, height, bg_color, text, text_color=None):
    """Create a small badge/pill"""
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = bg_color
    badge.line.fill.background()

    add_text_box(slide, left, top, width, height,
                 text, font_size=9, font_color=text_color or DARK_BG,
                 bold=True, alignment=PP_ALIGN.CENTER)
    return badge

def add_progress_bar(slide, left, top, width, height, progress, bg_color, fill_color):
    """Create a progress bar"""
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    # Fill
    fill_width = width * progress
    if fill_width > Inches(0.1):
        fill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, fill_width, height
        )
        fill.fill.solid()
        fill.fill.fore_color.rgb = fill_color
        fill.line.fill.background()

    return bg

def add_stat_number(slide, left, top, number, label, color):
    """Add a big stat number with label"""
    add_text_box(slide, left, top, Inches(1.5), Inches(0.5),
                 number, font_size=28, font_color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.4), Inches(1.5), Inches(0.3),
                 label, font_size=9, font_color=GRAY_LIGHT,
                 alignment=PP_ALIGN.CENTER)

def create_slide_1(prs):
    """Slide 1: Apertura y Re-encuadre"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # Logo REVIVE (texto estilizado)
    add_text_box(slide, Inches(3.5), Inches(1.5), Inches(6), Inches(1),
                 "RƎVIVE", font_size=72, font_color=EMERALD, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Subtítulo
    add_text_box(slide, Inches(2), Inches(2.6), Inches(9), Inches(0.6),
                 "Nueva Etapa", font_size=32, font_color=WHITE, bold=False,
                 alignment=PP_ALIGN.CENTER)

    # Mensaje clave
    add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(0.5),
                 '"Del experimento al activo comercial"',
                 font_size=24, font_color=VIOLET, bold=False,
                 alignment=PP_ALIGN.CENTER)

    # Timeline de evolución
    # Fase 1: Chatbot
    add_rounded_rectangle(slide, Inches(1.5), Inches(4.8), Inches(2.2), Inches(0.8), GRAY_DARK)
    add_text_box(slide, Inches(1.5), Inches(4.85), Inches(2.2), Inches(0.7),
                 "Chatbot", font_size=16, font_color=GRAY_LIGHT,
                 alignment=PP_ALIGN.CENTER)

    # Flecha 1
    add_arrow(slide, Inches(3.8), Inches(5), Inches(0.6), Inches(0.4), GRAY_DARK)

    # Fase 2: Aprendizaje
    add_rounded_rectangle(slide, Inches(4.5), Inches(4.8), Inches(2.2), Inches(0.8), GRAY_DARK)
    add_text_box(slide, Inches(4.5), Inches(4.85), Inches(2.2), Inches(0.7),
                 "Aprendizaje", font_size=16, font_color=GRAY_LIGHT,
                 alignment=PP_ALIGN.CENTER)

    # Flecha 2
    add_arrow(slide, Inches(6.8), Inches(5), Inches(0.6), Inches(0.4), EMERALD)

    # Fase 3: Plataforma (destacada)
    add_rounded_rectangle(slide, Inches(7.5), Inches(4.8), Inches(3.2), Inches(0.8), EMERALD)
    add_text_box(slide, Inches(7.5), Inches(4.85), Inches(3.2), Inches(0.7),
                 "Plataforma Profesional", font_size=16, font_color=DARK_BG, bold=True,
                 alignment=PP_ALIGN.CENTER)

def create_slide_2(prs):
    """Slide 2: Todo en Uno - Centralización de funcionalidades con flechas"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # Título
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "Todo lo que necesitas, en un solo lugar",
                 font_size=36, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Centro: Logo REVIVE grande
    center_x, center_y = 6.2, 3.5

    # Círculo central grande con glow
    add_circle(slide, Inches(center_x - 1.1), Inches(center_y - 1.1), Inches(2.2), EMERALD_DARK)
    add_circle(slide, Inches(center_x - 1), Inches(center_y - 1), Inches(2), EMERALD)
    add_text_box(slide, Inches(center_x - 1), Inches(center_y - 0.3), Inches(2), Inches(0.8),
                 "RƎVIVE", font_size=28, font_color=DARK_BG, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Funcionalidades con flechas hacia el centro
    features = [
        ("📅", "Agenda", BLUE, -2.8, -1.5, "right"),
        ("👥", "Clientes", VIOLET, 2.8, -1.5, "left"),
        ("💰", "Pagos", EMERALD_GLOW, -3.2, 0.8, "right"),
        ("📊", "Reportes", CYAN, 3.2, 0.8, "left"),
        ("📝", "Rutinas", PINK, -2, 2.2, "right"),
        ("💬", "Mensajes", BLUE_GLOW, 2, 2.2, "left"),
    ]

    for icon, name, color, offset_x, offset_y, arrow_dir in features:
        fx = center_x + offset_x
        fy = center_y + offset_y

        # Flecha hacia el centro
        if arrow_dir == "right":
            # Flecha apuntando a la derecha (hacia REVIVE)
            arrow_left = fx + 0.6
            arrow_width = abs(offset_x) - 1.6
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(arrow_left), Inches(fy - 0.15), Inches(arrow_width), Inches(0.3)
            )
        else:
            # Flecha apuntando a la izquierda (hacia REVIVE)
            arrow_left = center_x + 1
            arrow_width = abs(offset_x) - 1.6
            shape = slide.shapes.add_shape(
                MSO_SHAPE.LEFT_ARROW,
                Inches(arrow_left), Inches(fy - 0.15), Inches(arrow_width), Inches(0.3)
            )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        # Círculo de feature con glow
        add_circle(slide, Inches(fx - 0.55), Inches(fy - 0.55), Inches(1.1), GRAY_DARK)
        add_circle(slide, Inches(fx - 0.5), Inches(fy - 0.5), Inches(1), color)
        add_text_box(slide, Inches(fx - 0.5), Inches(fy - 0.35), Inches(1), Inches(0.5),
                     icon, font_size=24, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(fx - 0.7), Inches(fy + 0.5), Inches(1.4), Inches(0.4),
                     name, font_size=12, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

    # Mensaje inferior
    add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                 "Una plataforma diseñada para profesionales del bienestar",
                 font_size=18, font_color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)


def create_slide_2b(prs):
    """Slide 2b: Beneficios cuantitativos - Tu tiempo vale oro"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # Título positivo
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
                 "Lo que vas a ganar",
                 font_size=40, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(1), Inches(12), Inches(0.5),
                 "Resultados reales para tu negocio",
                 font_size=20, font_color=EMERALD, alignment=PP_ALIGN.CENTER)

    # 4 Cards de beneficios (incluyendo Nuevos Clientes Online)
    # Card 1: Tiempo recuperado
    add_glass_card(slide, Inches(0.5), Inches(1.7), Inches(2.9), Inches(2.2),
                   EMERALD, "Tiempo Recuperado", "⏱️")
    add_text_box(slide, Inches(0.5), Inches(2.3), Inches(2.9), Inches(0.7),
                 "5h+", font_size=48, font_color=EMERALD_GLOW, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(2.95), Inches(2.9), Inches(0.35),
                 "cada semana", font_size=16, font_color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(3.3), Inches(2.9), Inches(0.35),
                 "para tus clientes", font_size=11, font_color=GRAY_LIGHT,
                 alignment=PP_ALIGN.CENTER)

    # Card 2: Ingresos asegurados
    add_glass_card(slide, Inches(3.6), Inches(1.7), Inches(2.9), Inches(2.2),
                   BLUE, "Ingresos Asegurados", "💰")
    add_text_box(slide, Inches(3.6), Inches(2.3), Inches(2.9), Inches(0.7),
                 "100%", font_size=48, font_color=BLUE_GLOW, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.6), Inches(2.95), Inches(2.9), Inches(0.35),
                 "facturas controladas", font_size=16, font_color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.6), Inches(3.3), Inches(2.9), Inches(0.35),
                 "sin pagos olvidados", font_size=11, font_color=GRAY_LIGHT,
                 alignment=PP_ALIGN.CENTER)

    # Card 3: Clientes fidelizados
    add_glass_card(slide, Inches(6.7), Inches(1.7), Inches(2.9), Inches(2.2),
                   VIOLET, "Clientes Fidelizados", "🎯")
    add_text_box(slide, Inches(6.7), Inches(2.3), Inches(2.9), Inches(0.7),
                 "+40%", font_size=48, font_color=VIOLET_GLOW, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(6.7), Inches(2.95), Inches(2.9), Inches(0.35),
                 "adherencia", font_size=16, font_color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(6.7), Inches(3.3), Inches(2.9), Inches(0.35),
                 "seguimiento automático", font_size=11, font_color=GRAY_LIGHT,
                 alignment=PP_ALIGN.CENTER)

    # Card 4: NUEVOS CLIENTES ONLINE (destacada)
    add_glass_card(slide, Inches(9.8), Inches(1.7), Inches(2.9), Inches(2.2),
                   CYAN, "Nuevos Clientes", "🌐")
    add_text_box(slide, Inches(9.8), Inches(2.3), Inches(2.9), Inches(0.7),
                 "Online", font_size=36, font_color=CYAN, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(9.8), Inches(2.95), Inches(2.9), Inches(0.35),
                 "presencia digital", font_size=16, font_color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(9.8), Inches(3.3), Inches(2.9), Inches(0.35),
                 "capta clientes 24/7", font_size=11, font_color=GRAY_LIGHT,
                 alignment=PP_ALIGN.CENTER)

    # Fila inferior: Pills compactos (sin cajitas grandes)
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(12), Inches(0.4),
                 "Y además:", font_size=14, font_color=GRAY_LIGHT)

    pills = [
        ("📱 App profesional", EMERALD_DARK),
        ("🔔 Recordatorios auto", BLUE_DARK),
        ("📈 Reportes en vivo", VIOLET_DARK),
        ("🔒 Datos seguros", GRAY_DARK),
        ("📅 Agenda inteligente", EMERALD_DARK),
        ("💬 Mensajes directos", BLUE_DARK),
    ]

    # Primera fila de pills
    for i, (text, color) in enumerate(pills[:3]):
        left = 0.5 + (i * 4.2)
        add_rounded_rectangle(slide, Inches(left), Inches(4.6), Inches(3.8), Inches(0.55), color)
        add_text_box(slide, Inches(left), Inches(4.65), Inches(3.8), Inches(0.45),
                     text, font_size=14, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

    # Segunda fila de pills
    for i, (text, color) in enumerate(pills[3:]):
        left = 0.5 + (i * 4.2)
        add_rounded_rectangle(slide, Inches(left), Inches(5.3), Inches(3.8), Inches(0.55), color)
        add_text_box(slide, Inches(left), Inches(5.35), Inches(3.8), Inches(0.45),
                     text, font_size=14, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

    # Mensaje de cierre
    add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
                 "Menos gestión, más impacto",
                 font_size=22, font_color=EMERALD, bold=True, alignment=PP_ALIGN.CENTER)

def create_slide_3(prs):
    """Slide 3: La Solución - Teaser Visual con Glassmorphism Premium"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # Decorative background elements (subtle glows)
    # Top-left glow
    glow1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(4), Inches(4))
    glow1.fill.solid()
    glow1.fill.fore_color.rgb = RgbColor(0x16, 0x1B, 0x22)
    glow1.line.fill.background()

    # Bottom-right glow
    glow2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(5), Inches(5))
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = RgbColor(0x1A, 0x16, 0x22)
    glow2.line.fill.background()

    # Título con gradiente simulado
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 "Tu ventaja competitiva",
                 font_size=38, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════════
    # CARD 1: AGENDA (Top-Left) - Emerald theme
    # ═══════════════════════════════════════════════════════════════
    add_glass_card(slide, Inches(0.5), Inches(1.1), Inches(3.3), Inches(2.6),
                   EMERALD, "Agenda", "📅")

    # Session cards inside
    add_mini_card(slide, Inches(0.65), Inches(1.7), Inches(3), Inches(0.75))
    # Avatar circle with gradient effect
    add_circle(slide, Inches(0.75), Inches(1.8), Inches(0.55), BLUE)
    add_text_box(slide, Inches(0.78), Inches(1.88), Inches(0.5), Inches(0.35),
                 "MG", font_size=11, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.4), Inches(1.78), Inches(2), Inches(0.3),
                 "María García", font_size=12, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(1.4), Inches(2.02), Inches(1.2), Inches(0.25),
                 "10:00 - 11:00", font_size=10, font_color=EMERALD_GLOW)
    add_badge(slide, Inches(2.7), Inches(2.02), Inches(0.8), Inches(0.25),
              EMERALD, "Personal", WHITE)

    add_mini_card(slide, Inches(0.65), Inches(2.55), Inches(3), Inches(0.75))
    add_circle(slide, Inches(0.75), Inches(2.65), Inches(0.55), VIOLET)
    add_text_box(slide, Inches(0.78), Inches(2.73), Inches(0.5), Inches(0.35),
                 "CR", font_size=11, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.4), Inches(2.63), Inches(2), Inches(0.3),
                 "Carlos Ruiz", font_size=12, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(1.4), Inches(2.87), Inches(1.2), Inches(0.25),
                 "11:30 - 12:30", font_size=10, font_color=VIOLET_GLOW)
    add_badge(slide, Inches(2.7), Inches(2.87), Inches(0.8), Inches(0.25),
              VIOLET, "Grupal", WHITE)

    # Stats row
    add_text_box(slide, Inches(0.7), Inches(3.4), Inches(1), Inches(0.25),
                 "Hoy: 6 sesiones", font_size=9, font_color=GRAY_LIGHT)
    add_text_box(slide, Inches(2.2), Inches(3.4), Inches(1.4), Inches(0.25),
                 "Semana: 24/28", font_size=9, font_color=EMERALD_GLOW)

    # ═══════════════════════════════════════════════════════════════
    # CARD 2: CLIENTES (Top-Right) - Blue theme
    # ═══════════════════════════════════════════════════════════════
    add_glass_card(slide, Inches(4), Inches(1.1), Inches(3.3), Inches(2.6),
                   BLUE, "Clientes", "👥")

    # Client rows
    add_mini_card(slide, Inches(4.15), Inches(1.7), Inches(3), Inches(0.65))
    add_text_box(slide, Inches(4.25), Inches(1.78), Inches(1.5), Inches(0.3),
                 "Ana López", font_size=11, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(4.25), Inches(1.98), Inches(1.2), Inches(0.25),
                 "12 sesiones", font_size=9, font_color=GRAY_LIGHT)
    add_badge(slide, Inches(6.25), Inches(1.82), Inches(0.7), Inches(0.35),
              EMERALD, "95%", WHITE)

    add_mini_card(slide, Inches(4.15), Inches(2.45), Inches(3), Inches(0.65))
    add_text_box(slide, Inches(4.25), Inches(2.53), Inches(1.5), Inches(0.3),
                 "Pedro Martín", font_size=11, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(4.25), Inches(2.73), Inches(1.2), Inches(0.25),
                 "8 sesiones", font_size=9, font_color=GRAY_LIGHT)
    add_badge(slide, Inches(6.25), Inches(2.57), Inches(0.7), Inches(0.35),
              RgbColor(0xFB, 0xBF, 0x24), "72%", DARK_BG)

    # Alert badge
    add_mini_card(slide, Inches(4.15), Inches(3.2), Inches(3), Inches(0.4),
                  RgbColor(0x1E, 0x14, 0x14))
    add_text_box(slide, Inches(4.25), Inches(3.22), Inches(2.8), Inches(0.35),
                 "⚠️  2 clientes necesitan seguimiento",
                 font_size=9, font_color=RgbColor(0xFB, 0x92, 0x3C))

    # ═══════════════════════════════════════════════════════════════
    # CARD 3: FACTURACIÓN (Bottom-Left) - Violet theme
    # ═══════════════════════════════════════════════════════════════
    add_glass_card(slide, Inches(0.5), Inches(3.9), Inches(3.3), Inches(2.6),
                   VIOLET, "Facturación", "💰")

    # Big number
    add_text_box(slide, Inches(0.65), Inches(4.5), Inches(3), Inches(0.6),
                 "$2.450.000", font_size=32, font_color=EMERALD_GLOW, bold=True)
    add_text_box(slide, Inches(0.65), Inches(5), Inches(1.5), Inches(0.25),
                 "Ingresos Enero", font_size=10, font_color=GRAY_LIGHT)
    add_badge(slide, Inches(2.2), Inches(5), Inches(1), Inches(0.25),
              EMERALD_DARK, "+12% ↑", EMERALD_GLOW)

    # Mini stats
    add_mini_card(slide, Inches(0.65), Inches(5.4), Inches(1.4), Inches(0.7))
    add_text_box(slide, Inches(0.7), Inches(5.45), Inches(1.3), Inches(0.25),
                 "Cobrado", font_size=8, font_color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.7), Inches(5.7), Inches(1.3), Inches(0.35),
                 "85%", font_size=18, font_color=EMERALD, bold=True, alignment=PP_ALIGN.CENTER)

    add_mini_card(slide, Inches(2.15), Inches(5.4), Inches(1.5), Inches(0.7),
                  RgbColor(0x22, 0x16, 0x16))
    add_text_box(slide, Inches(2.2), Inches(5.45), Inches(1.4), Inches(0.25),
                 "Pendiente", font_size=8, font_color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.2), Inches(5.7), Inches(1.4), Inches(0.35),
                 "3", font_size=18, font_color=RgbColor(0xFB, 0x71, 0x85), bold=True, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════════
    # CARD 4: REPORTES (Bottom-Right) - Cyan theme
    # ═══════════════════════════════════════════════════════════════
    add_glass_card(slide, Inches(4), Inches(3.9), Inches(3.3), Inches(2.6),
                   CYAN, "Reportes", "📊")

    # Adherence section
    add_text_box(slide, Inches(4.15), Inches(4.5), Inches(2), Inches(0.25),
                 "Adherencia promedio", font_size=10, font_color=GRAY_LIGHT)

    # Progress bar
    add_progress_bar(slide, Inches(4.15), Inches(4.8), Inches(3), Inches(0.25),
                     0.78, RgbColor(0x1E, 0x1E, 0x28), EMERALD)
    add_text_box(slide, Inches(4.15), Inches(5.1), Inches(1), Inches(0.25),
                 "78%", font_size=14, font_color=EMERALD, bold=True)
    add_text_box(slide, Inches(5), Inches(5.12), Inches(2), Inches(0.25),
                 "+5% vs mes anterior", font_size=9, font_color=EMERALD_GLOW)

    # Mini ranking
    add_mini_card(slide, Inches(4.15), Inches(5.5), Inches(3), Inches(0.9))
    add_text_box(slide, Inches(4.25), Inches(5.55), Inches(2), Inches(0.25),
                 "🏆 Top Adherencia", font_size=9, font_color=RgbColor(0xFB, 0xBF, 0x24), bold=True)
    add_text_box(slide, Inches(4.25), Inches(5.85), Inches(2.8), Inches(0.25),
                 "1. María G. (95%)  2. Ana L. (92%)  3. Luis M. (89%)",
                 font_size=8, font_color=GRAY_LIGHT)

    # ═══════════════════════════════════════════════════════════════
    # RIGHT SIDE: Message and CTA
    # ═══════════════════════════════════════════════════════════════

    # Main message with styled background
    msg_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.6), Inches(2), Inches(5.2), Inches(2.2)
    )
    msg_bg.fill.solid()
    msg_bg.fill.fore_color.rgb = RgbColor(0x18, 0x18, 0x24)
    msg_bg.line.color.rgb = EMERALD_DARK
    msg_bg.line.width = Pt(2)

    add_text_box(slide, Inches(7.9), Inches(2.3), Inches(4.8), Inches(0.8),
                 "Una plataforma.", font_size=36, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(7.9), Inches(2.9), Inches(4.8), Inches(0.8),
                 "Todo tu negocio.", font_size=36, font_color=EMERALD_GLOW, bold=True)

    # Feature pills
    pills_data = [
        ("Agenda", EMERALD, 7.9),
        ("Clientes", BLUE, 9.3),
        ("Pagos", VIOLET, 10.55),
        ("Reportes", CYAN, 11.6)
    ]
    for text, color, left in pills_data:
        add_badge(slide, Inches(left), Inches(3.9), Inches(1.1), Inches(0.35),
                  color, text, WHITE)

    # CTA transition
    add_text_box(slide, Inches(7.9), Inches(5.5), Inches(5), Inches(0.5),
                 "Pero mejor te lo enseño en vivo...",
                 font_size=20, font_color=VIOLET_GLOW, bold=True)

    # Decorative arrow
    add_arrow(slide, Inches(11.5), Inches(5.55), Inches(0.8), Inches(0.4), VIOLET)

def create_slide_4(prs):
    """Slide 4: De Friends & Family a Business Partners"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # Título
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
                 "Una nueva relación",
                 font_size=40, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Dos columnas: Antes vs Ahora
    # Columna izquierda: ANTES (Friends & Family)
    add_rounded_rectangle(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.2),
                          RgbColor(0x1A, 0x1A, 0x24), line_color=GRAY_DARK)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
                 "Hasta ahora", font_size=14, font_color=GRAY_LIGHT,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(0.6),
                 "Friends & Family", font_size=28, font_color=GRAY_LIGHT, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Items del antes
    antes_items = [
        ("🎓", "Aprendizaje mutuo"),
        ("🤝", "Colaboración informal"),
        ("💡", "Experimentación"),
        ("🎁", "Sin compromisos"),
    ]
    for i, (icon, text) in enumerate(antes_items):
        y_pos = 2.9 + (i * 0.6)
        add_text_box(slide, Inches(1.5), Inches(y_pos), Inches(4.5), Inches(0.5),
                     f"{icon}  {text}", font_size=16, font_color=GRAY_MEDIUM)

    # Flecha central grande
    add_arrow(slide, Inches(6.1), Inches(3.2), Inches(1.2), Inches(0.8), EMERALD)

    # Columna derecha: AHORA (Business Partners)
    add_rounded_rectangle(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(4.2),
                          GRAY_DARK, line_color=EMERALD)
    add_text_box(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.5),
                 "Nueva etapa", font_size=14, font_color=EMERALD,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(0.6),
                 "Business Partners", font_size=28, font_color=EMERALD, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Items del ahora
    ahora_items = [
        ("🚀", "Producto profesional"),
        ("📈", "Crecimiento conjunto"),
        ("💼", "Relación comercial"),
        ("🎯", "Objetivos compartidos"),
    ]
    for i, (icon, text) in enumerate(ahora_items):
        y_pos = 2.9 + (i * 0.6)
        add_text_box(slide, Inches(7.7), Inches(y_pos), Inches(4.5), Inches(0.5),
                     f"{icon}  {text}", font_size=16, font_color=WHITE)

    # Mensaje de transición
    add_rounded_rectangle(slide, Inches(2.5), Inches(5.9), Inches(8), Inches(1.2),
                          RgbColor(0x16, 0x20, 0x16), line_color=EMERALD_DARK)
    add_text_box(slide, Inches(2.5), Inches(6), Inches(8), Inches(0.5),
                 "El siguiente paso natural",
                 font_size=14, font_color=EMERALD, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.5), Inches(6.4), Inches(8), Inches(0.5),
                 "Profesionalizar nuestra colaboración para crecer juntos",
                 font_size=18, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

def create_slide_5(prs):
    """Slide 5: Parking Lot con ideas y espacios vacíos"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # Título
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
                 "Roadmap de Ideas",
                 font_size=40, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtítulo
    add_text_box(slide, Inches(0.5), Inches(1), Inches(12), Inches(0.5),
                 "Genial para la versión 2.0 - Lo anotamos aquí",
                 font_size=18, font_color=VIOLET, alignment=PP_ALIGN.CENTER)

    # Ideas ya definidas (cards con contenido)
    ideas_existentes = [
        ("🤖", "IA Coaching", "Asistente inteligente\npara recomendaciones", EMERALD),
        ("📱", "App Móvil Nativa", "iOS y Android\ncon notificaciones push", BLUE),
        ("🔗", "Integraciones", "Google Calendar,\nStripe, WhatsApp API", VIOLET),
    ]

    for i, (icon, title, desc, color) in enumerate(ideas_existentes):
        left = 0.7 + (i * 4.2)
        # Card con idea
        add_rounded_rectangle(slide, Inches(left), Inches(1.7), Inches(3.8), Inches(1.8),
                              GRAY_DARK, line_color=color)
        add_text_box(slide, Inches(left), Inches(1.8), Inches(3.8), Inches(0.5),
                     icon, font_size=28, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(left), Inches(2.3), Inches(3.8), Inches(0.4),
                     title, font_size=16, font_color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(left), Inches(2.7), Inches(3.8), Inches(0.7),
                     desc, font_size=11, font_color=GRAY_LIGHT,
                     alignment=PP_ALIGN.CENTER)

    # Espacios vacíos para anotar en vivo (con icono de +)
    add_text_box(slide, Inches(0.5), Inches(3.8), Inches(12), Inches(0.4),
                 "Ideas de Christian:", font_size=14, font_color=GRAY_LIGHT)

    for i in range(3):
        left = 0.7 + (i * 4.2)
        # Card vacía con borde punteado simulado
        add_rounded_rectangle(slide, Inches(left), Inches(4.3), Inches(3.8), Inches(1.5),
                              RgbColor(0x1A, 0x1A, 0x24), line_color=GRAY_DARK)
        add_text_box(slide, Inches(left), Inches(4.7), Inches(3.8), Inches(0.6),
                     "+", font_size=36, font_color=GRAY_DARK,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(left), Inches(5.3), Inches(3.8), Inches(0.4),
                     "Tu idea aquí", font_size=12, font_color=GRAY_DARK,
                     alignment=PP_ALIGN.CENTER)

    # Mensaje inferior
    add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
                 "Priorizamos juntos según el impacto en tu negocio",
                 font_size=16, font_color=EMERALD, alignment=PP_ALIGN.CENTER)

def create_slide_6(prs):
    """Slide 6: Cierre - Oportunidad de crecer juntos"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # Fondo decorativo: Gráfico de crecimiento ascendente simulado
    # Barras de crecimiento
    bar_data = [
        (2.5, 5.5, 0.8),   # Barra 1 (más pequeña)
        (3.5, 4.8, 1.5),   # Barra 2
        (4.5, 4.0, 2.3),   # Barra 3
        (5.5, 3.2, 3.1),   # Barra 4
        (6.5, 2.4, 3.9),   # Barra 5 (más grande)
    ]
    for left, top, height in bar_data:
        add_rounded_rectangle(slide, Inches(left), Inches(top), Inches(0.7), Inches(height),
                              RgbColor(0x1A, 0x2A, 0x1A))

    # Línea de tendencia (flecha diagonal)
    add_arrow(slide, Inches(2.3), Inches(5.5), Inches(5.5), Inches(0.5), EMERALD_DARK)

    # Logo REVIVE
    add_text_box(slide, Inches(3.5), Inches(0.8), Inches(6), Inches(0.9),
                 "RƎVIVE", font_size=56, font_color=EMERALD, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Mensaje principal
    add_text_box(slide, Inches(1.5), Inches(1.9), Inches(10), Inches(0.8),
                 "Crezcamos juntos",
                 font_size=44, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtítulo
    add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.6),
                 "La oportunidad de construir algo grande",
                 font_size=22, font_color=EMERALD_GLOW, alignment=PP_ALIGN.CENTER)

    # Cards de propuesta de valor
    # Card 1: Tú aportas
    add_rounded_rectangle(slide, Inches(1), Inches(3.6), Inches(3.5), Inches(2.2),
                          GRAY_DARK, line_color=BLUE)
    add_text_box(slide, Inches(1), Inches(3.75), Inches(3.5), Inches(0.4),
                 "Tú aportas", font_size=14, font_color=BLUE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.2), Inches(4.2), Inches(3.1), Inches(1.4),
                 "👤 Tu experiencia\n📣 Tu feedback real\n🎯 Tu visión de mercado",
                 font_size=14, font_color=WHITE)

    # Símbolo de unión
    add_circle(slide, Inches(4.9), Inches(4.2), Inches(0.9), EMERALD)
    add_text_box(slide, Inches(4.9), Inches(4.35), Inches(0.9), Inches(0.6),
                 "+", font_size=36, font_color=DARK_BG, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Card 2: Nosotros aportamos
    add_rounded_rectangle(slide, Inches(6.2), Inches(3.6), Inches(3.5), Inches(2.2),
                          GRAY_DARK, line_color=VIOLET)
    add_text_box(slide, Inches(6.2), Inches(3.75), Inches(3.5), Inches(0.4),
                 "Nosotros aportamos", font_size=14, font_color=VIOLET,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(6.4), Inches(4.2), Inches(3.1), Inches(1.4),
                 "💻 Tecnología de punta\n🚀 Desarrollo continuo\n🛡️ Soporte dedicado",
                 font_size=14, font_color=WHITE)

    # Card 3: Resultado
    add_rounded_rectangle(slide, Inches(10), Inches(3.6), Inches(2.5), Inches(2.2),
                          EMERALD_DARK, line_color=EMERALD)
    add_text_box(slide, Inches(10), Inches(3.75), Inches(2.5), Inches(0.4),
                 "Resultado", font_size=14, font_color=EMERALD_GLOW,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(10), Inches(4.3), Inches(2.5), Inches(0.8),
                 "📈", font_size=48, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(10), Inches(5.1), Inches(2.5), Inches(0.5),
                 "Éxito\ncompartido", font_size=14, font_color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Mensaje de cierre
    add_rounded_rectangle(slide, Inches(3), Inches(6.1), Inches(7), Inches(0.9),
                          RgbColor(0x16, 0x16, 0x22), line_color=EMERALD)
    add_text_box(slide, Inches(3), Inches(6.25), Inches(7), Inches(0.6),
                 "¿Construimos el futuro de REVIVE juntos?",
                 font_size=20, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

def main():
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Create all slides
    create_slide_1(prs)      # Apertura: REVIVE Nueva Etapa
    create_slide_2(prs)      # Todo en uno: Centralización
    create_slide_2b(prs)     # Lo que vas a ganar: Beneficios
    create_slide_3(prs)      # Teaser Visual: Glassmorphism
    create_slide_4(prs)      # Friends & Family → Business Partners
    create_slide_5(prs)      # Roadmap de Ideas
    create_slide_6(prs)      # Crezcamos juntos

    # Save presentation
    output_path = '/Users/juan.poblete/Documents/Personal/TheMindFactory/TheMindFactory/REVIVE_Reunion_Christian.pptx'
    prs.save(output_path)
    print(f"Presentación guardada en: {output_path}")
    print(f"Total de slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
